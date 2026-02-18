"""
Convert a PowerPoint deck (.pptx) into a single structured text file.

This script is designed for internal physical therapy documentation:
- Accuracy > speed
- Readability > compression
- Output is optimized for later RAG chunking (but chunking is NOT implemented here).

Usage
-----
    python pptx_to_structured_text.py path/to/deck.pptx
    python pptx_to_structured_text.py path/to/deck.pptx -o path/to/deck.txt

OCR note
--------
This script uses `pytesseract`, which requires the external Tesseract binary.
If Tesseract is not installed or not on PATH, extraction still proceeds and a
warning is written to a separate `<pptx_stem>_RAG_errors.txt` report.

Image description note
----------------------
Image descriptions can come from:
- PowerPoint metadata (alt text / title), and/or
- A local captioning model via `transformers` (optional).

If captioning dependencies are missing, extraction still succeeds with a clear
warning in the same `<pptx_stem>_RAG_errors.txt` report.

Preserved elements
------------------
- Slide boundaries
- Text boxes (in approximate reading order)
- Bullet hierarchy (indentation levels)
- Speaker notes (if present)
- Image RAG context (computer-vision description + OCR text)
- Optional raw OCR text from embedded slide images, linked to the slide it came from

Non-goals (intentionally NOT implemented)
----------------------------------------
- Chunking
- Embeddings / vectorization
- RAG retrieval
- Scoring, LLM prompting, or any external services
"""

from __future__ import annotations

import argparse
import io
import os
import re
import sys
from datetime import datetime
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import List, Optional, Sequence, Tuple


def _require_imports() -> tuple[object, object, object]:
    """
    Import required third-party libraries with a clear error message if missing.

    Returns
    -------
    (pptx, PIL.Image, pytesseract)
    """

    try:
        import pptx  # type: ignore
        from PIL import Image  # type: ignore
        import pytesseract  # type: ignore
    except ModuleNotFoundError as exc:
        missing = str(exc).split("'")[-2] if "'" in str(exc) else str(exc)
        raise RuntimeError(
            "Missing required dependency. Install with:\n"
            "  pip install python-pptx pillow pytesseract\n"
            f"Missing module: {missing}"
        ) from exc

    return pptx, Image, pytesseract


def load_pptx(pptx_path: Path):
    """Load a .pptx as a python-pptx Presentation."""

    pptx, _, _ = _require_imports()
    from pptx import Presentation  # type: ignore

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")
    if pptx_path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file: {pptx_path}")
    return Presentation(str(pptx_path))


@dataclass(frozen=True)
class TextRun:
    """
    A text line extracted from a slide.

    Attributes
    ----------
    text:
        The text content (already stripped/normalized).
    indent_level:
        Bullet indentation level (0+). This is PowerPoint paragraph level.
    is_bullet:
        Whether the paragraph is represented as a bullet in output.
    source_label:
        Optional label for where this text came from (e.g., shape name).
    """

    text: str
    indent_level: int
    is_bullet: bool
    source_label: Optional[str] = None


@dataclass(frozen=True)
class SlideContent:
    slide_number: int
    title: str
    text_lines: List[TextRun]
    bullet_lines: List[TextRun]
    speaker_notes: str
    image_description_sections: List[Tuple[str, str]]  # (label, description)
    image_rag_sections: List[Tuple[str, str]]  # (label, rag_context)
    image_ocr_sections: List[Tuple[str, str]]  # (label, ocr_text)


@dataclass(frozen=True)
class ExtractedImage:
    label: str
    image_bytes: bytes
    metadata_description: str


@dataclass(frozen=True)
class IngestWarning:
    slide_number: int
    warning_type: str
    image_label: str
    message: str


def _emu(value) -> int:
    """Convert pptx length-like values to int EMU; treat missing as huge for sorting."""

    if value is None:
        return 10**18
    try:
        return int(value)
    except Exception:
        return 10**18


def _shape_sort_key(shape) -> tuple[int, int, int]:
    # Approximate reading order: top-to-bottom, then left-to-right, then stable id.
    return (_emu(getattr(shape, "top", None)), _emu(getattr(shape, "left", None)), int(getattr(shape, "shape_id", 0)))


def _normalize_line(text: str) -> str:
    return " ".join(text.replace("\r", "\n").split()).strip()


def _is_bulleted_paragraph(paragraph) -> bool:
    """
    Best-effort bullet detection using python-pptx's underlying XML.

    PowerPoint supports bullets at level 0, so paragraph.level alone is not enough.
    """

    # Public hint: indented paragraphs are usually bullets.
    try:
        if int(getattr(paragraph, "level", 0)) > 0:
            return True
    except Exception:
        pass

    # XML hint: <a:buChar>, <a:buAutoNum>, <a:buBlip>, and absence of <a:buNone>.
    ppr = getattr(paragraph, "_pPr", None)
    if ppr is None:
        return False

    try:
        from pptx.oxml.ns import qn  # type: ignore
    except Exception:
        return False

    if ppr.find(qn("a:buNone")) is not None:
        return False
    if ppr.find(qn("a:buChar")) is not None:
        return True
    if ppr.find(qn("a:buAutoNum")) is not None:
        return True
    if ppr.find(qn("a:buBlip")) is not None:
        return True
    return False


def extract_text_shapes(slide) -> tuple[str, List[TextRun], List[TextRun]]:
    """
    Extract title text plus other text (bulleted and non-bulleted) from a slide.

    Returns
    -------
    (title, text_lines, bullet_lines)
    """

    title_text = ""
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title_text = _normalize_line(title_shape.text_frame.text or "")

    text_lines: List[TextRun] = []
    bullet_lines: List[TextRun] = []

    shapes = sorted(list(slide.shapes), key=_shape_sort_key)
    for shape in shapes:
        if shape is title_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text_frame = shape.text_frame
        shape_label = getattr(shape, "name", None)

        for paragraph in text_frame.paragraphs:
            line = _normalize_line(getattr(paragraph, "text", "") or "")
            if not line:
                continue
            indent_level = int(getattr(paragraph, "level", 0) or 0)
            is_bullet = _is_bulleted_paragraph(paragraph)
            run = TextRun(text=line, indent_level=max(0, indent_level), is_bullet=is_bullet, source_label=shape_label)
            if is_bullet:
                bullet_lines.append(run)
            else:
                text_lines.append(run)

    return title_text, text_lines, bullet_lines


def extract_speaker_notes(slide) -> str:
    """Extract speaker notes from a slide, if present."""

    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""

    notes_text_frame = getattr(notes_slide, "notes_text_frame", None)
    if notes_text_frame is None:
        return ""

    text = notes_text_frame.text or ""
    # PowerPoint often includes a default placeholder line; keep only meaningful text.
    normalized = "\n".join(
        line.strip()
        for line in (text.replace("\r\n", "\n").replace("\r", "\n")).split("\n")
        if line.strip()
    ).strip()
    return normalized


def _extract_shape_alt_text(shape) -> str:
    """
    Extract alt text / title metadata from a shape when available.

    Returns normalized text or empty string.
    """

    direct_candidates = [
        getattr(shape, "alternative_text", None),
        getattr(shape, "alt_text", None),
        getattr(shape, "title", None),
    ]
    for value in direct_candidates:
        if isinstance(value, str):
            normalized = _normalize_line(value)
            if normalized:
                return normalized

    element = getattr(shape, "_element", None)
    if element is None:
        return ""

    # python-pptx exposes the underlying XML; cNvPr usually carries `descr` and `title`.
    for node in element.iter():
        tag = getattr(node, "tag", "")
        if isinstance(tag, str) and tag.endswith("}cNvPr"):
            descr = _normalize_line((node.get("descr") or ""))
            title = _normalize_line((node.get("title") or ""))
            if descr:
                return descr
            if title:
                return title
    return ""


def extract_images(slide) -> List[ExtractedImage]:
    """
    Extract embedded images from a slide.

    Returns
    -------
    List of extracted image records for each image found in the slide.
    """

    images: List[ExtractedImage] = []

    # Import constants lazily (only available when python-pptx is installed).
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception:
        MSO_SHAPE_TYPE = None

    for idx, shape in enumerate(slide.shapes, start=1):
        shape_type = getattr(shape, "shape_type", None)
        has_image = hasattr(shape, "image")
        if MSO_SHAPE_TYPE is not None:
            if shape_type != MSO_SHAPE_TYPE.PICTURE and not has_image:
                continue
        else:
            if not has_image:
                continue

        try:
            image = shape.image
            blob = image.blob  # bytes
            filename = getattr(image, "filename", None) or f"image_{idx}"
            label = f"image {idx} ({filename})"
            metadata_description = _extract_shape_alt_text(shape)
            images.append(
                ExtractedImage(
                    label=label,
                    image_bytes=blob,
                    metadata_description=metadata_description,
                )
            )
        except Exception:
            # Best-effort: ignore shapes that can't be extracted.
            continue

    return images


def _check_ocr_available(pytesseract_module, *, tesseract_cmd: Optional[str] = None) -> Optional[str]:
    """
    Returns an error message if OCR cannot run, else None.

    This checks for the external Tesseract binary, which pytesseract requires.
    """

    if tesseract_cmd:
        try:
            pytesseract_module.pytesseract.tesseract_cmd = tesseract_cmd
        except Exception:
            # Keep best-effort behavior and fall through to version check.
            pass

    try:
        pytesseract_module.get_tesseract_version()
        return None
    except Exception as exc:
        location_hint = f" (configured path: {tesseract_cmd})" if tesseract_cmd else ""
        return (
            "OCR unavailable: Tesseract engine not found or not runnable. "
            "Install Tesseract and ensure it is on PATH, or pass --tesseract-cmd / set TESSERACT_CMD."
            f"{location_hint} "
            f"Details: {exc}"
        )


def _normalize_easyocr_languages(languages: str) -> tuple[str, ...]:
    """
    Parse EasyOCR language input into a stable tuple.
    """

    normalized = (languages or "").replace(";", ",").strip()
    if not normalized:
        return ("en",)

    mapped: List[str] = []
    for token in normalized.split(","):
        lang = token.strip().lower()
        if not lang:
            continue
        # Common bridge from Tesseract code.
        if lang == "eng":
            lang = "en"
        mapped.append(lang)

    return tuple(mapped) if mapped else ("en",)


@lru_cache(maxsize=8)
def _load_easyocr_reader(language_codes: tuple[str, ...], gpu: bool):
    """
    Lazy-load and cache an EasyOCR Reader.
    """

    try:
        import easyocr  # type: ignore
    except ModuleNotFoundError as exc:
        missing = str(exc).split("'")[-2] if "'" in str(exc) else str(exc)
        return None, f"OCR unavailable (easyocr): missing dependency '{missing}'. Install with: pip install easyocr"

    try:
        reader = easyocr.Reader(
            list(language_codes),
            gpu=gpu,
            verbose=False,
            download_enabled=False,
        )
        return reader, None
    except TypeError:
        # Older EasyOCR versions may not accept verbose= or download_enabled=.
        try:
            reader = easyocr.Reader(
                list(language_codes),
                gpu=gpu,
                download_enabled=False,
            )
            return reader, None
        except TypeError:
            return (
                None,
                "OCR unavailable (easyocr): this EasyOCR version cannot enforce offline mode "
                "(missing download_enabled support). Upgrade EasyOCR.",
            )
        except Exception as exc:
            return None, f"OCR unavailable (easyocr): failed to initialize reader. Details: {exc}"
    except Exception as exc:
        return None, f"OCR unavailable (easyocr): failed to initialize reader. Details: {exc}"


def run_ocr_easyocr(
    image_bytes: bytes,
    *,
    languages: str = "en",
    gpu: bool = False,
) -> tuple[str, Optional[str]]:
    """
    Run OCR using EasyOCR.
    """

    _, Image, _ = _require_imports()

    language_codes = _normalize_easyocr_languages(languages)
    reader, reader_err = _load_easyocr_reader(language_codes, gpu)
    if reader_err is not None or reader is None:
        return "", reader_err or "OCR unavailable (easyocr)."

    try:
        import numpy as np  # type: ignore
    except ModuleNotFoundError as exc:
        missing = str(exc).split("'")[-2] if "'" in str(exc) else str(exc)
        return "", f"OCR unavailable (easyocr): missing dependency '{missing}'. Install with: pip install easyocr"

    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            image_array = np.array(img.convert("RGB"))
            raw = reader.readtext(image_array, detail=0, paragraph=True)
    except Exception as exc:
        return "", f"OCR failed while processing image with easyocr: {exc}"

    lines = [str(item).strip() for item in (raw or []) if str(item).strip()]
    normalized = "\n".join(lines).strip()
    return normalized, None


def run_ocr(
    image_bytes: bytes,
    *,
    backend: str = "tesseract",
    language: str = "eng",
    tesseract_cmd: Optional[str] = None,
    easyocr_languages: str = "en",
    easyocr_gpu: bool = False,
) -> tuple[str, Optional[str]]:
    """
    Run OCR on an image blob.

    Returns
    -------
    (ocr_text, error_message)

    Notes
    -----
    - If OCR is unavailable (e.g., missing Tesseract), returns ("", error_message).
    - Image decoding and OCR are intentionally conservative; accuracy > speed.
    """

    if backend == "easyocr":
        return run_ocr_easyocr(
            image_bytes,
            languages=easyocr_languages,
            gpu=easyocr_gpu,
        )

    _, Image, pytesseract = _require_imports()

    ocr_error = _check_ocr_available(pytesseract, tesseract_cmd=tesseract_cmd)
    if ocr_error is not None:
        return "", ocr_error

    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            # Mild normalization for OCR robustness.
            img = img.convert("RGB")
            text = pytesseract.image_to_string(img, lang=language) or ""
            normalized = "\n".join(line.rstrip() for line in text.splitlines()).strip()
            return normalized, None
    except Exception as exc:
        return "", f"OCR failed while processing image: {exc}"


@lru_cache(maxsize=4)
def _load_caption_pipeline(model_name: str):
    """
    Lazy-load an image captioning pipeline.

    Returns
    -------
    (pipeline_obj, error_message)
    """

    try:
        import torch  # type: ignore
        from transformers import pipeline  # type: ignore
    except ModuleNotFoundError as exc:
        missing = str(exc).split("'")[-2] if "'" in str(exc) else str(exc)
        return (
            None,
            "Image description unavailable: missing captioning dependency "
            f"'{missing}'. Install with: pip install transformers torch",
        )

    try:
        device = 0 if torch.cuda.is_available() else -1
    except Exception:
        device = -1

    # Hard offline guard for SCIF-safe usage: never attempt Hub calls at runtime.
    os.environ["HF_HUB_OFFLINE"] = "1"
    os.environ["TRANSFORMERS_OFFLINE"] = "1"

    last_error: Optional[Exception] = None
    for task in ("image-to-text", "image-text-to-text"):
        try:
            cap = pipeline(task, model=model_name, device=device, local_files_only=True)
            return cap, None
        except Exception as exc:
            last_error = exc
            continue
    return (
        None,
        "Image description unavailable: could not load caption model "
        f"'{model_name}' from local files in offline mode. Details: {last_error}",
    )


def run_image_caption(
    image_bytes: bytes,
    *,
    model_name: str,
) -> tuple[str, Optional[str]]:
    """
    Generate a visual description for an image blob.

    Returns
    -------
    (description, error_message)
    """

    _, Image, _ = _require_imports()

    captioner, captioner_err = _load_caption_pipeline(model_name)
    if captioner_err is not None or captioner is None:
        return "", captioner_err or "Image description unavailable."

    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            img = img.convert("RGB")
            task_name = str(getattr(captioner, "task", "") or "").strip().lower()
            if task_name == "image-text-to-text":
                raw = captioner(images=img, text="Describe this image.")
            else:
                raw = captioner(img)
    except Exception as exc:
        return "", f"Image description failed while processing image: {exc}"

    text = ""
    if isinstance(raw, list) and raw:
        top = raw[0]
        if isinstance(top, dict):
            text = str(top.get("generated_text", "") or "").strip()
        else:
            text = str(top).strip()
    elif isinstance(raw, dict):
        text = str(raw.get("generated_text", "") or "").strip()
    else:
        text = str(raw).strip()

    normalized = _normalize_line(text)
    lower = normalized.lower()
    for prefix in ("describe this image.", "describe this image:"):
        if lower.startswith(prefix):
            normalized = normalized[len(prefix) :].strip(" -:\t")
            break
    if not normalized:
        return "", "Image description failed: caption model returned empty output."
    return normalized, None


def build_image_description(
    image: ExtractedImage,
    *,
    mode: str,
    caption_model: str,
) -> tuple[str, Optional[str]]:
    """
    Build a textual image description according to the selected mode.
    """

    metadata = image.metadata_description.strip()

    if mode == "none":
        return "", None
    if mode == "metadata":
        return metadata, None
    if mode == "caption":
        caption, cap_err = run_image_caption(image.image_bytes, model_name=caption_model)
        return (caption, None) if cap_err is None else ("", cap_err)

    # auto mode: prefer metadata when author-provided alt text exists, else caption.
    if metadata:
        return metadata, None
    caption, cap_err = run_image_caption(image.image_bytes, model_name=caption_model)
    return (caption, None) if cap_err is None else ("", cap_err)


def _summarize_ocr_for_rag(
    ocr_text: str,
    *,
    max_lines: int = 0,
    max_chars: int = 0,
) -> str:
    """
    Convert raw OCR output into a compact, retrieval-friendly snippet.
    """

    cleaned_lines: List[str] = []
    for raw in (ocr_text or "").splitlines():
        normalized = _normalize_line(raw)
        if normalized:
            cleaned_lines.append(normalized)

    if not cleaned_lines:
        return ""

    if max_lines and max_lines > 0:
        kept = cleaned_lines[: max(1, max_lines)]
    else:
        kept = cleaned_lines
    joined = " | ".join(kept).strip()
    if max_chars and max_chars > 0 and len(joined) > max_chars:
        joined = joined[: max_chars - 1].rstrip() + "..."
    return joined


def _is_weak_visual_description(text: str) -> bool:
    """
    Heuristic filter for low-information captions like "s" or "icon".
    """

    normalized = _normalize_line(text or "")
    if not normalized:
        return True

    lowered = normalized.lower()
    weak_literals = {"s", "icon", "logo", "image", "graphic", "photo", "com"}
    if lowered in weak_literals:
        return True

    tokens = [t for t in re.split(r"\s+", normalized) if t]
    if len(tokens) < 2:
        return True
    if len(normalized) < 10:
        return True
    return False


def _infer_visual_from_ocr(ocr_summary: str) -> str:
    """
    Build a coarse visual description using OCR-derived cues.
    """

    low = (ocr_summary or "").lower()
    has_resolution = bool(re.search(r"\b\d{3,4}\s*[xX]\s*\d{3,4}\b", ocr_summary or ""))
    has_grid = "grid" in low
    has_test = "test" in low
    has_numeric = bool(re.search(r"\b\d+(?:\.\d+)?\b", ocr_summary or ""))
    has_bars = "color" in low or "colour" in low or "bar" in low

    if has_resolution or has_grid or has_test:
        return "Text-heavy test-pattern/reference graphic."
    if has_numeric and has_bars:
        return "Text-heavy technical graphic with numeric/color labels."
    if has_numeric:
        return "Text-heavy technical visual with numeric labels."
    return "Text-heavy visual with embedded labels."


def build_image_rag_context(
    *,
    visual_description: str,
    ocr_text: str,
    rag_ocr_max_lines: int = 0,
    rag_ocr_max_chars: int = 0,
) -> str:
    """
    Build a RAG-oriented image context from CV + OCR.
    """

    desc = _normalize_line(visual_description or "")
    ocr_summary = _summarize_ocr_for_rag(
        ocr_text,
        max_lines=rag_ocr_max_lines,
        max_chars=rag_ocr_max_chars,
    )

    parts: List[str] = []
    if desc and not _is_weak_visual_description(desc):
        parts.append(f"Visual: {desc}")
    elif ocr_summary:
        parts.append(f"Visual: {_infer_visual_from_ocr(ocr_summary)}")
    elif desc:
        parts.append(f"Visual: {desc}")
    if ocr_summary:
        parts.append(f"Detected text: {ocr_summary}")

    if not parts:
        return "No usable visual description or OCR text extracted."
    return "\n".join(parts)


def _render_bullets(bullets: Sequence[TextRun]) -> str:
    out_lines: List[str] = []
    for run in bullets:
        indent = "  " * max(0, int(run.indent_level))
        out_lines.append(f"{indent}- {run.text}")
    return "\n".join(out_lines).rstrip()


def _render_text_lines(text_lines: Sequence[TextRun]) -> str:
    # Preserve shape boundaries by emitting a label when the source changes.
    out_lines: List[str] = []
    last_label: Optional[str] = None
    for run in text_lines:
        label = run.source_label or None
        if label and label != last_label:
            out_lines.append(f"[Text Box: {label}]")
            last_label = label
        out_lines.append(run.text)
    return "\n".join(out_lines).rstrip()


def _render_slide(content: SlideContent) -> str:
    lines: List[str] = []
    lines.append(f"=== Slide {content.slide_number} ===")
    lines.append("Title:")
    lines.append(content.title or "")
    lines.append("")
    lines.append("Text:")
    lines.append(_render_text_lines(content.text_lines) or "")
    lines.append("")
    lines.append("Bullets:")
    lines.append(_render_bullets(content.bullet_lines) or "")
    has_image_context = bool(content.image_rag_sections or content.image_ocr_sections)
    if has_image_context:
        lines.append("")
        lines.append("Image RAG Context:")
        if content.image_rag_sections:
            for label, rag_text in content.image_rag_sections:
                lines.append(f"[RAG context for {label}]")
                lines.append(rag_text.strip() if rag_text.strip() else "")
                lines.append("")
            if lines and lines[-1] == "":
                lines.pop()
        else:
            lines.append("")
    if content.image_ocr_sections:
        lines.append("")
        lines.append("Image Text (Raw OCR):")
        for label, ocr_text in content.image_ocr_sections:
            lines.append(f"[OCR from {label}]")
            lines.append(ocr_text.strip() if ocr_text.strip() else "")
            lines.append("")
        if lines and lines[-1] == "":
            lines.pop()
    if content.speaker_notes.strip():
        lines.append("")
        lines.append("Speaker Notes:")
        lines.append(content.speaker_notes.strip())
    return "\n".join(lines).rstrip() + "\n"


def write_output(slides: Sequence[SlideContent], output_path: Path) -> None:
    """Write extracted slide content into a single structured text file."""

    output_path.parent.mkdir(parents=True, exist_ok=True)
    text = "\n".join(_render_slide(s) for s in slides).rstrip() + "\n"
    output_path.write_text(text, encoding="utf-8", errors="replace")


def write_warning_report(
    warnings: Sequence[IngestWarning],
    errors_path: Path,
    *,
    source_pptx: Path,
) -> None:
    """Write ingest warnings to a dedicated report file."""

    errors_path.parent.mkdir(parents=True, exist_ok=True)

    lines: List[str] = []
    lines.append("=== RAG Ingest Warnings ===")
    lines.append(f"Source PPTX: {source_pptx}")
    lines.append(f"Generated: {datetime.now().isoformat(timespec='seconds')}")
    lines.append(f"Total warnings: {len(warnings)}")
    lines.append("")

    if not warnings:
        lines.append("No warnings detected.")
    else:
        for entry in warnings:
            lines.append(f"[Slide {entry.slide_number}] [{entry.warning_type}] {entry.image_label}")
            lines.append(entry.message.strip())
            lines.append("")
        if lines and lines[-1] == "":
            lines.pop()

    errors_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8", errors="replace")


def convert_pptx_to_text(
    pptx_path: Path,
    *,
    output_path: Optional[Path] = None,
    errors_path: Optional[Path] = None,
    ocr_backend: str = "easyocr",
    ocr_language: str = "eng",
    tesseract_cmd: Optional[str] = None,
    easyocr_languages: str = "en",
    easyocr_gpu: bool = False,
    image_description_mode: str = "caption",
    caption_model: str = "Salesforce/blip-image-captioning-base",
    rag_ocr_max_lines: int = 0,
    rag_ocr_max_chars: int = 0,
    include_raw_image_ocr: bool = False,
) -> Path:
    """
    Convert a PPTX into a single structured text output file.

    Parameters
    ----------
    pptx_path:
        Path to the input .pptx.
    output_path:
        Optional output file path. Defaults to `<pptx_stem>_extracted.txt`
        next to the input file.
    errors_path:
        Optional warning report path. Defaults to `<pptx_stem>_RAG_errors.txt`
        next to the input file.
    ocr_backend:
        OCR engine to use: "tesseract" or "easyocr".
    ocr_language:
        Tesseract language code (default: "eng"). Ignored for easyocr backend.
    tesseract_cmd:
        Optional explicit path to the Tesseract executable. If omitted, this
        function uses the `TESSERACT_CMD` environment variable if present.
        Ignored for easyocr backend.
    easyocr_languages:
        EasyOCR language codes as comma-separated values (default: "en").
    easyocr_gpu:
        If True, EasyOCR uses GPU when available.
    image_description_mode:
        Image description strategy: "none", "metadata", "caption", or "auto".
    caption_model:
        Hugging Face model id for local captioning when caption mode is used.
    rag_ocr_max_lines:
        Max OCR lines to include in `Image RAG Context` per image (0 = no limit).
    rag_ocr_max_chars:
        Max OCR characters to include in `Image RAG Context` per image (0 = no limit).
    include_raw_image_ocr:
        If True, include a separate raw OCR section in the output.
    """

    presentation = load_pptx(pptx_path)
    if output_path is None:
        output_path = pptx_path.with_name(f"{pptx_path.stem}_extracted.txt")
    if errors_path is None:
        errors_path = pptx_path.with_name(f"{pptx_path.stem}_RAG_errors.txt")
    if ocr_backend == "tesseract" and tesseract_cmd is None:
        env_tesseract_cmd = (os.environ.get("TESSERACT_CMD") or "").strip()
        if env_tesseract_cmd:
            tesseract_cmd = env_tesseract_cmd

    warnings: List[IngestWarning] = []
    slide_contents: List[SlideContent] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        title, text_lines, bullet_lines = extract_text_shapes(slide)
        speaker_notes = extract_speaker_notes(slide)

        image_descriptions: List[Tuple[str, str]] = []
        image_rag_sections: List[Tuple[str, str]] = []
        image_sections: List[Tuple[str, str]] = []
        images = extract_images(slide)
        for image in images:
            description, description_err = build_image_description(
                image,
                mode=image_description_mode,
                caption_model=caption_model,
            )
            image_descriptions.append((image.label, description))
            if description_err is not None:
                warnings.append(
                    IngestWarning(
                        slide_number=idx,
                        warning_type="IMAGE_DESCRIPTION",
                        image_label=image.label,
                        message=description_err,
                    )
                )

            ocr_text, ocr_err = run_ocr(
                image.image_bytes,
                backend=ocr_backend,
                language=ocr_language,
                tesseract_cmd=tesseract_cmd,
                easyocr_languages=easyocr_languages,
                easyocr_gpu=easyocr_gpu,
            )
            if ocr_err is not None:
                warnings.append(
                    IngestWarning(
                        slide_number=idx,
                        warning_type="OCR",
                        image_label=image.label,
                        message=ocr_err,
                    )
                )
            rag_context = build_image_rag_context(
                visual_description=description,
                ocr_text=ocr_text,
                rag_ocr_max_lines=rag_ocr_max_lines,
                rag_ocr_max_chars=rag_ocr_max_chars,
            )
            image_rag_sections.append((image.label, rag_context))

            if include_raw_image_ocr:
                image_sections.append((image.label, ocr_text))

        slide_contents.append(
            SlideContent(
                slide_number=idx,
                title=title,
                text_lines=text_lines,
                bullet_lines=bullet_lines,
                speaker_notes=speaker_notes,
                image_description_sections=image_descriptions,
                image_rag_sections=image_rag_sections,
                image_ocr_sections=image_sections,
            )
        )

    write_output(slide_contents, output_path)
    write_warning_report(warnings, errors_path, source_pptx=pptx_path)
    return output_path


def _parse_args(argv: Optional[Sequence[str]] = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert a .pptx into structured text with OCR and image descriptions.")
    parser.add_argument("pptx_path", type=str, help="Path to the input .pptx file")
    parser.add_argument(
        "-o",
        "--output",
        type=str,
        default=None,
        help="Optional output .txt path (default: next to the .pptx)",
    )
    parser.add_argument(
        "--ocr-backend",
        choices=("tesseract", "easyocr"),
        default="easyocr",
        help="OCR backend (default: easyocr)",
    )
    parser.add_argument(
        "--ocr-language",
        type=str,
        default="eng",
        help="Tesseract language code (default: eng; ignored for easyocr)",
    )
    parser.add_argument(
        "--tesseract-cmd",
        type=str,
        default=None,
        help="Optional path to tesseract executable (alternative to PATH or TESSERACT_CMD; tesseract backend only)",
    )
    parser.add_argument(
        "--easyocr-languages",
        type=str,
        default="en",
        help="EasyOCR language codes, comma-separated (default: en)",
    )
    parser.add_argument(
        "--easyocr-gpu",
        action="store_true",
        help="Enable GPU for EasyOCR when available",
    )
    parser.add_argument(
        "--image-description-mode",
        choices=("none", "metadata", "caption", "auto"),
        default="caption",
        help="How to generate image descriptions (default: caption)",
    )
    parser.add_argument(
        "--caption-model",
        type=str,
        default="Salesforce/blip-image-captioning-base",
        help="Caption model for --image-description-mode caption/auto (default: Salesforce/blip-image-captioning-base)",
    )
    parser.add_argument(
        "--errors-output",
        type=str,
        default=None,
        help="Optional warning report path (default: <pptx_stem>_RAG_errors.txt)",
    )
    parser.add_argument(
        "--rag-ocr-max-lines",
        type=int,
        default=0,
        help="Max OCR lines in Image RAG Context per image (0 = no limit)",
    )
    parser.add_argument(
        "--rag-ocr-max-chars",
        type=int,
        default=0,
        help="Max OCR chars in Image RAG Context per image (0 = no limit)",
    )
    parser.add_argument(
        "--include-raw-image-ocr",
        action="store_true",
        help="Include a separate raw OCR section in the extracted text",
    )
    return parser.parse_args(argv)


def main(argv: Optional[Sequence[str]] = None) -> int:
    try:
        args = _parse_args(argv)
        pptx_path = Path(args.pptx_path).expanduser().resolve()
        output_path = Path(args.output).expanduser().resolve() if args.output else None
        errors_path = Path(args.errors_output).expanduser().resolve() if args.errors_output else None
        out = convert_pptx_to_text(
            pptx_path,
            output_path=output_path,
            errors_path=errors_path,
            ocr_backend=args.ocr_backend,
            ocr_language=args.ocr_language,
            tesseract_cmd=args.tesseract_cmd,
            easyocr_languages=args.easyocr_languages,
            easyocr_gpu=args.easyocr_gpu,
            image_description_mode=args.image_description_mode,
            caption_model=args.caption_model,
            rag_ocr_max_lines=args.rag_ocr_max_lines,
            rag_ocr_max_chars=args.rag_ocr_max_chars,
            include_raw_image_ocr=args.include_raw_image_ocr,
        )
        # Minimal, useful CLI output (not verbose tracing).
        print(str(out))
        if errors_path is None:
            errors_path = pptx_path.with_name(f"{pptx_path.stem}_RAG_errors.txt")
        print(f"Warnings report: {errors_path}")
        return 0
    except RuntimeError as exc:
        print(str(exc), file=sys.stderr)
        return 2
    except Exception as exc:
        print(f"Failed: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
